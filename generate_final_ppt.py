from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

# Colors from previous scripts for consistency
BRAND_COLOR = RGBColor(0, 71, 64)  # Dark teal from add_team_names.py
ACCENT_GREEN = RGBColor(52, 211, 153) # Matches dashboard positive green

def update_text_in_shape(shape, text, font_size=None, bold=None, color=None):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.text = text
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            if font_size: run.font.size = Pt(font_size)
            if bold is not None: run.font.bold = bold
            if color: run.font.color.rgb = color

def generate_pitch_deck():
    template_path = 'd:/Projects/Hackathon/TCD-GCGC/Idea Sprint 3.0 PPT Template.pptx'
    output_path = 'd:/Projects/Hackathon/TCD-GCGC/Final_AI_Simulator_Pitch_Deck.pptx'
    
    prs = Presentation(template_path)
    
    # --- Slide 0: Title & Problem Code ---
    slide0 = prs.slides[0]
    team_names = "Team: Krishna Srinivas, Yashwanth, Saravana Sai, Venkat"
    project_title = "AI PRODUCT STRATEGY SIMULATOR [AI04]"
    
    for shape in slide0.shapes:
        if shape.has_text_frame:
            if "PROBLEM STATEMENT CODE" in shape.text:
                update_text_in_shape(shape, project_title, font_size=36, bold=True, color=BRAND_COLOR)
            if "KRISHNA SR" in shape.text:
                update_text_in_shape(shape, team_names, font_size=24, color=BRAND_COLOR)

    # --- Slide 1-8: Content Overwrite ---
    # We will clear and add content to the main body areas
    content = [
        ("THE PROBLEM", [
            "Companies lose millions on bad pricing and product pivots.",
            "Market forces (Competitors, Customers, Investors) are hard to predict.",
            "Current strategy testing is slow, manual, and relies on stale data.",
            "Teams lack a sandbox to 'play through' market reactions before launch."
        ]),
        ("OUR SOLUTION", [
            "Autonomous Multi-Agent Market Simulator.",
            "Converts static PRDs into dynamic, six-stage strategy simulations.",
            "Self-correcting AI agents (ADK) simulate real-world reactions.",
            "Integrated OSINT pulls live world events (News/Reddit) for high-fidelity grounding."
        ]),
        ("ENGINEERING ARCHITECTURE", [
            "Frontend: React + ECharts 'War Room' dashboard.",
            "Backend: FastAPI with a transparent 6-stage pipeline architecture.",
            "Intelligence: Gemini-driven PRD extraction and Multi-Agent decision loops.",
            "Execution: ADK Agents for reasoning + Deterministic fallbacks for speed/predictability."
        ]),
        ("UNIQUE INNOVATIONS", [
            "PRD-to-Sim: Zero-config simulation from any PDF/Doc upload.",
            "Auditable Reasoning: Every agent decision is traceable and visible in UI.",
            "Crisis Injection: Manual override to shock the market with custom scenarios.",
            "Scenario Recommendation: AI suggests the best pricing pivot based on simulation history."
        ]),
        ("TECHNICAL DEPTH", [
            "Agent To User Interface (A2UI): Transparent mapping of latent AI steps.",
            "Resilience: Dual-mode parsing (LLM + Structural) ensures 100% uptime.",
            "OSINT Grounding: Real-time RSS feeds prevent 'hallucination' in market signals.",
            "Performance: Parallelized agent loops with state-tracked progression."
        ]),
        ("DASHBOARD: THE WAR ROOM", [
            "Probability of Success (PoS) Gauge: Instant confidence readout.",
            "Market Arena: Dynamic supply/demand curves and revenue timelines.",
            "Sentiment Heatmap: Geographic and demographic sentiment breakdown.",
            "Agent Activity: Live, per-step reasoning logs from 4 distinct personas."
        ]),
        ("BUSINESS IMPACT", [
            "Accelerated GTM: Strategy testing in minutes, not months.",
            "Risk Mitigation: Pre-simulate competitor price wars and investor panic.",
            "Better Margins: Recommendation engine finds the revenue sweet spot.",
            "Explainable Strategy: Board-room ready visualizations of AI predictions."
        ]),
        ("THE FUTURE: SCALE & ROADMAP", [
            "Deep Personalities: Simulated high-fidelity buyer personas.",
            "Scenario Libraries: Pre-built templates for recessions, tech booms, or regulations.",
            "Enterprise API: Strategy-as-a-Service for existing PM tools.",
            "Collaborative War Rooms: Multiplayer strategy co-pilot for teams."
        ]),
    ]

    for i, (title, points) in enumerate(content):
        if (i + 1) < len(prs.slides):
            slide = prs.slides[i + 1]
            # Try to find title and body placeholders
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if "Training" in shape.text or "Slide" in shape.text or not shape.text:
                        # Clear old template text and insert ours
                        tf = shape.text_frame
                        tf.clear()
                        p = tf.paragraphs[0]
                        p.text = title
                        p.font.bold = True
                        p.font.size = Pt(28)
                        for point in points:
                            p = tf.add_paragraph()
                            p.text = "• " + point
                            p.font.size = Pt(18)
                            p.space_after = Pt(10)

    # --- Slide 9: Thank You ---
    slide9 = prs.slides[9]
    for shape in slide9.shapes:
        if shape.has_text_frame and "Thank" in shape.text:
            update_text_in_shape(shape, "THANK YOU!", font_size=44, bold=True, color=BRAND_COLOR)

    prs.save(output_path)
    print(f"✅ Final deck created at: {output_path}")

if __name__ == '__main__':
    generate_pitch_deck()
