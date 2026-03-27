from pptx import Presentation

def analyze_template(path):
    try:
        prs = Presentation(path)
        print(f"Template Analysis for: {path}")
        print("-" * 30)
        for i, layout in enumerate(prs.slide_layouts):
            print(f"Layout {i}: {layout.name}")
            for placeholder in layout.placeholders:
                print(f"  - [{placeholder.placeholder_format.idx}] {placeholder.name} (type: {placeholder.placeholder_format.type})")
        
        print("\nExisting Slides:")
        for i, slide in enumerate(prs.slides):
            print(f"Slide {i}:")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    print(f"  - Shape: {shape.name} | Text: {shape.text[:50]}...")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == '__main__':
    analyze_template('d:/Projects/Hackathon/TCD-GCGC/Idea Sprint 3.0 PPT Template.pptx')
