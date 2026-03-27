from pptx import Presentation

def analyze_template_slides(path):
    try:
        prs = Presentation(path)
        print(f"Total Slides: {len(prs.slides)}")
        for i, slide in enumerate(prs.slides):
            print(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip().replace('\n', ' ')
                    if text:
                        print(f"  Shape: {shape.name} | Text: {text[:100]}")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == '__main__':
    analyze_template_slides('d:/Projects/Hackathon/TCD-GCGC/Idea Sprint 3.0 PPT Template.pptx')
