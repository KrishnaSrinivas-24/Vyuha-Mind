from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Load current deck
prs = Presentation('AI_Product_Strategy_Simulator_Pitch_Deck.pptx')
slide = prs.slides[0]

# Template formatting for team names
FONT_SIZE = 24  # pt
FONT_COLOR = RGBColor(0x00, 0x47, 0x40)  # 004740 - dark teal
ALIGNMENT = PP_ALIGN.CENTER

# Add text box for team names
team_names = ["Krishna Srinivas", "Yashwanth", "Saravana Sai", "Venkat"]
left = Emu(-78)
top = Emu(3898100)
width = Emu(12192000)
height = Emu(1791219)

text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

for i, name in enumerate(team_names):
    if i > 0:
        p = text_frame.add_paragraph()
    else:
        p = text_frame.paragraphs[0]
    
    p.text = name
    p.alignment = ALIGNMENT
    p.level = 0
    
    # Format the run
    run = p.runs[0]
    run.font.size = Pt(FONT_SIZE)
    run.font.color.rgb = FONT_COLOR

# Save the updated deck
prs.save('AI_Product_Strategy_Simulator_Pitch_Deck.pptx')
print("✅ Team names added to first slide!")
print(f"Names added: {', '.join(team_names)}")
