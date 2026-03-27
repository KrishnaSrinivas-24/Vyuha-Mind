from pptx import Presentation

p = Presentation('AI_Product_Strategy_Simulator_Pitch_Deck.pptx')
slide = p.slides[0]
print('First slide now has', len(slide.shapes), 'shapes')
for i, shape in enumerate(slide.shapes):
    if hasattr(shape, "text"):
        print(f'  Shape {i}: {shape.text[:60]}')
    else:
        print(f'  Shape {i}: (non-text shape)')
