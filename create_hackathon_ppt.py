from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TITLE_COLOR = RGBColor(21, 34, 56)
BODY_COLOR = RGBColor(48, 64, 89)
ACCENT_COLOR = RGBColor(28, 126, 214)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    title_tf = slide.shapes.title.text_frame
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = TITLE_COLOR

    sub_tf = slide.placeholders[1].text_frame
    sub_tf.paragraphs[0].font.size = Pt(20)
    sub_tf.paragraphs[0].font.color.rgb = ACCENT_COLOR


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    title_tf = slide.shapes.title.text_frame
    title_tf.paragraphs[0].font.size = Pt(34)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = TITLE_COLOR

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, item in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(9)


def add_demo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "DEMONSTRATION"
    title_tf = slide.shapes.title.text_frame
    title_tf.paragraphs[0].font.size = Pt(34)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = TITLE_COLOR

    # Left block
    left = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.9), Inches(4.9))
    ltf = left.text_frame
    ltf.word_wrap = True
    ltf.clear()

    left_points = [
        "Functional Prototype (Live):",
        "1) Upload PRD (PDF/DOC/DOCX/TXT/MD)",
        "2) Auto-extract product, audience, pricing, scenario",
        "3) Launch multi-agent simulation",
        "4) Observe live dashboards + KPI movement",
        "5) Inspect full agent pipeline outputs",
        "6) Review recommendation engine strategy"
    ]

    for i, item in enumerate(left_points):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(21 if i == 0 else 19)
        p.font.bold = i == 0
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(8)

    # Right block placeholder
    right = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(6.0), Inches(4.9))
    rtf = right.text_frame
    rtf.word_wrap = True
    rtf.clear()

    rp0 = rtf.paragraphs[0]
    rp0.text = "Demo Visual Placeholder"
    rp0.font.size = Pt(22)
    rp0.font.bold = True
    rp0.font.color.rgb = ACCENT_COLOR

    rp1 = rtf.add_paragraph()
    rp1.text = "Insert screenshot of War Room dashboard or use live run during pitch."
    rp1.font.size = Pt(18)
    rp1.font.color.rgb = BODY_COLOR

    rp2 = rtf.add_paragraph()
    rp2.text = "Critical outputs shown: Input Handler, Orchestrator, Market Analyzer, Simulation Loop, Evaluation, Recommendation."
    rp2.font.size = Pt(17)
    rp2.font.color.rgb = BODY_COLOR


def build_deck(output_path):
    prs = Presentation()

    add_title_slide(
        prs,
        "Autonomous Multi-Agent Product Strategy Simulator",
        "Hackathon Pitch Deck"
    )

    add_bullet_slide(prs, "PROBLEM STATEMENT", [
        "Teams launching products in volatile markets lack fast, evidence-based strategy testing.",
        "Target users: founders, product managers, strategy teams, innovation labs.",
        "Current methods are slow, manual, and depend on fragmented market intelligence.",
        "Most tools do not model customer, competitor, and investor behavior together.",
        "Result: weak pricing decisions, delayed pivots, and increased go-to-market risk."
    ])

    add_bullet_slide(prs, "OUR SOLUTION", [
        "We built a multi-agent strategy simulator that turns product inputs into actionable strategy recommendations.",
        "Users can upload a PRD, auto-extract required fields, and run a full simulation instantly.",
        "The platform simulates customer, competitor, and investor responses under dynamic market scenarios.",
        "A live dashboard shows KPI shifts, agent reasoning, risk, and recommended pricing paths.",
        "Value proposition: faster decisions, lower strategy risk, and clearer market-fit confidence."
    ])

    add_bullet_slide(prs, "TECHNICAL APPROACH", [
        "Frontend: React + Vite dashboard with real-time simulation playback and pipeline visibility.",
        "Backend: FastAPI orchestration engine with structured endpoints (/pipeline/run, /prd/parse).",
        "AI/Intelligence: Gemini-based PRD extraction + fallback parser for resilience.",
        "Simulation workflow: Input Handler -> Orchestrator -> Market Analyzer -> Multi-Agent Loop -> Evaluation -> Recommendation.",
        "Engineering focus: deterministic fallbacks, CORS-enabled integration, and test-backed reliability."
    ])

    add_bullet_slide(prs, "KEY FEATURES", [
        "PRD Upload & Auto-Extraction: parse PDF/DOC/DOCX/TXT/MD into simulation-ready fields.",
        "Live Multi-Agent Simulation: customer, competitor, and investor reasoning over sequential steps.",
        "Interactive War Room: moving charts, KPI panel, sentiment heatmap, and competitor matrix.",
        "Full Pipeline Transparency: six engine-stage outputs shown directly in frontend.",
        "Action Controls: play/pause/step/speed, scenario injection, and recommendation-driven decisions."
    ])

    add_bullet_slide(prs, "INNOVATION & IMPACT", [
        "Innovation: combines PRD intelligence + multi-agent market simulation in one practical workflow.",
        "Differentiator: exposes every engine stage, making AI strategy recommendations auditable.",
        "Effectiveness: converts uncertain strategy choices into measurable success/risk signals.",
        "Application areas: SaaS pricing, consumer product launches, GTM planning, crisis strategy response.",
        "Real-world impact: improves strategic agility and helps teams pivot before market losses escalate."
    ])

    add_demo_slide(prs)

    add_bullet_slide(prs, "CHALLENGES & FUTURE SCOPE", [
        "Challenge: model/API availability and quota limits during real-time parsing and analysis.",
        "Mitigation: robust fallback parsers and deterministic simulation paths to avoid downtime.",
        "Challenge: keeping UI responsive while simulating multi-step agent behavior.",
        "Mitigation: optimized state progression and controlled playback loops.",
        "Future scope: deeper scenario libraries, richer persona models, collaboration workflows, and enterprise reporting."
    ])

    add_bullet_slide(prs, "CONCLUSION", [
        "We addressed a high-impact problem: uncertain strategy decisions in dynamic markets.",
        "Our solution combines AI extraction, multi-agent simulation, and transparent decision support.",
        "The system is practical, explainable, and already running with a working frontend-backend pipeline.",
        "This approach can help teams make faster, safer, and more data-backed product decisions.",
        "Closing statement: from guesswork to guided strategy, in one simulation loop."
    ])

    prs.save(output_path)


if __name__ == '__main__':
    build_deck('d:/Projects/Hackathon/TCD-GCGC/AI_Product_Strategy_Simulator_Pitch_Deck.pptx')
    print('created')
